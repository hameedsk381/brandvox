import io
from uuid import UUID
from datetime import datetime
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from sqlalchemy.ext.asyncio import AsyncSession
from typing import Optional, Dict, Any
from app.services.analytics_service import get_dashboard_data
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

async def generate_excel_report(
    db: AsyncSession,
    agency_id: Optional[UUID] = None,
    client_id: Optional[UUID] = None,
    location_id: Optional[UUID] = None,
    report_type: str = "monthly"
) -> bytes:
    # 1. Fetch dashboard metrics
    data = await get_dashboard_data(db, agency_id, client_id, location_id)
    
    # 2. Build workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Reputation Summary"
    
    # Enable grid lines
    ws.views.sheetView[0].showGridLines = True
    
    # Styles
    title_font = Font(name="Calibri", size=16, bold=True, color="FFFFFF")
    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    bold_font = Font(name="Calibri", size=11, bold=True)
    regular_font = Font(name="Calibri", size=11)
    
    title_fill = PatternFill(start_color="4F46E5", end_color="4F46E5", fill_type="solid") # Indigo 600
    header_fill = PatternFill(start_color="312E81", end_color="312E81", fill_type="solid") # Indigo 900
    stripe_fill = PatternFill(start_color="F3F4F6", end_color="F3F4F6", fill_type="solid") # Gray 100
    
    border_side = Side(border_style="thin", color="D1D5DB")
    thin_border = Border(left=border_side, right=border_side, top=border_side, bottom=border_side)
    
    # Title Row
    ws.merge_cells("A1:D1")
    ws["A1"] = f"ReputationOS AI - {report_type.capitalize()} Brand Intelligence Report"
    ws["A1"].font = title_font
    ws["A1"].fill = title_fill
    ws["A1"].alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 40
    
    # Space
    ws.append([])
    
    # Key Metrics Header
    ws.append(["Metric", "Value", "Target", "Status"])
    ws.row_dimensions[3].height = 24
    for col in range(1, 5):
        cell = ws.cell(row=3, column=col)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(vertical="center")
        cell.border = thin_border
        
    # KPI Rows
    metrics = [
        ("Reputation Score", data["reputation_score"], "85.0+", "Healthy" if data["reputation_score"] >= 80 else "Needs Attention"),
        ("Average Rating", data["avg_rating"], "4.5★", "Good" if data["avg_rating"] >= 4.0 else "Action Required"),
        ("Total Reviews", data["total_reviews"], "-", "Active"),
        ("Response Rate", f"{int(data['response_rate'] * 100)}%", "90%+", "Optimal" if data["response_rate"] >= 0.85 else "Low response"),
        ("Sentiment Index", f"{int(data['sentiment_score'] * 100)}%", "80%+", "Satisfactory" if data["sentiment_score"] >= 0.7 else "Volatile")
    ]
    
    for row_idx, (metric, val, target, status) in enumerate(metrics, start=4):
        ws.append([metric, val, target, status])
        ws.row_dimensions[row_idx].height = 20
        for col in range(1, 5):
            cell = ws.cell(row=row_idx, column=col)
            cell.font = regular_font
            cell.border = thin_border
            if row_idx % 2 == 1:
                cell.fill = stripe_fill
                
    # Auto-adjust column widths
    from openpyxl.cell.cell import MergedCell
    for col in ws.columns:
        non_merged = [c for c in col if not isinstance(c, MergedCell)]
        if non_merged:
            max_len = max(len(str(cell.value or '')) for cell in non_merged)
            ws.column_dimensions[non_merged[0].column_letter].width = max(max_len + 3, 12)
        
    # Save to memory stream
    stream = io.BytesIO()
    wb.save(stream)
    return stream.getvalue()

async def generate_pdf_report(
    db: AsyncSession,
    agency_id: Optional[UUID] = None,
    client_id: Optional[UUID] = None,
    location_id: Optional[UUID] = None,
    report_type: str = "monthly"
) -> bytes:
    # 1. Fetch dashboard metrics
    data = await get_dashboard_data(db, agency_id, client_id, location_id)
    
    # 2. Setup document
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=letter,
        rightMargin=36,
        leftMargin=36,
        topMargin=36,
        bottomMargin=36
    )
    
    # Styles
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "ReportTitle",
        parent=styles["Title"],
        fontName="Helvetica-Bold",
        fontSize=24,
        textColor=colors.HexColor("#4F46E5"),
        spaceAfter=15
    )
    subtitle_style = ParagraphStyle(
        "ReportSubtitle",
        parent=styles["Normal"],
        fontSize=12,
        textColor=colors.HexColor("#4B5563"),
        spaceAfter=30
    )
    h2_style = ParagraphStyle(
        "SectionHeader",
        parent=styles["Heading2"],
        fontSize=16,
        textColor=colors.HexColor("#1E1B4B"),
        spaceBefore=20,
        spaceAfter=10
    )
    body_style = ParagraphStyle(
        "BodyTextCustom",
        parent=styles["Normal"],
        fontSize=10,
        textColor=colors.HexColor("#374151"),
        spaceAfter=10
    )
    
    story = []
    
    # Header
    story.append(Paragraph(f"ReputationOS AI — {report_type.capitalize()} Brand Audit", title_style))
    story.append(Paragraph(f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')} | Target: Platform Analytics Summary", subtitle_style))
    story.append(Spacer(1, 10))
    
    # Summary Paragraph
    story.append(Paragraph("Executive Overview", h2_style))
    overview_text = (
        f"During this performance cycle, your reputation score is registered at <b>{data['reputation_score']}/100</b>, "
        f"with an average customer rating of <b>{data['avg_rating']} stars</b> based on a total volume of <b>{data['total_reviews']} reviews</b>. "
        f"The system logged a reply rate of <b>{int(data['response_rate'] * 100)}%</b> for incoming reviews. "
        f"Sentiment classification maps to <b>{int(data['sentiment_score'] * 100)}% positive</b> reception."
    )
    story.append(Paragraph(overview_text, body_style))
    story.append(Spacer(1, 15))
    
    # Metrics Table
    story.append(Paragraph("Key Performance Diagnostics", h2_style))
    table_data = [
        ["Metric Breakdown", "Current Performance Value", "Industry Benchmark Status"],
        ["Reputation Index", f"{data['reputation_score']}/100", "Top 15% Tier" if data["reputation_score"] >= 80 else "Median Tier"],
        ["Star Rating", f"{data['avg_rating']}★", "Optimal" if data["avg_rating"] >= 4.2 else "Review Campaign Suggested"],
        ["Reply Coverage", f"{int(data['response_rate'] * 100)}%", "Satisfies SLA" if data["response_rate"] >= 0.85 else "Action Needed"],
        ["Positive Sentiment Ratio", f"{int(data['sentiment_score'] * 100)}%", "Highly Positive" if data["sentiment_score"] >= 0.75 else "Volatile"]
    ]
    
    t = Table(table_data, colWidths=[200, 160, 180])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor("#312E81")),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 10),
        ('TOPPADDING', (0, 0), (-1, 0), 10),
        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor("#F9FAFB")]),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor("#E5E7EB")),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
        ('BOTTOMPADDING', (0, 1), (-1, -1), 8),
        ('TOPPADDING', (0, 1), (-1, -1), 8),
    ]))
    story.append(t)
    
    # Build Document
    doc.build(story)
    return buffer.getvalue()


async def generate_pptx_report(
    db: AsyncSession,
    agency_id: Optional[UUID] = None,
    client_id: Optional[UUID] = None,
    location_id: Optional[UUID] = None,
    report_type: str = "monthly"
) -> bytes:
    # 1. Fetch dashboard metrics
    data = await get_dashboard_data(db, agency_id, client_id, location_id)
    
    prs = Presentation()
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Use blank slide layout
    blank_layout = prs.slide_layouts[6]
    
    # Define premium color scheme (slate-dark theme)
    bg_color = RGBColor(15, 23, 42)       # Slate 900
    card_bg_color = RGBColor(30, 41, 59)  # Slate 800
    primary_color = RGBColor(99, 102, 241) # Indigo 500
    text_white = RGBColor(248, 250, 252)  # Slate 50
    text_muted = RGBColor(148, 163, 184)  # Slate 400
    
    # Helper to set slide background color
    def set_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    # Top decorative brand bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.4))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = primary_color
    top_bar.line.fill.background()
    
    # Title Text Frame
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "ReputationOS AI"
    p.font.name = "Arial"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = primary_color
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = f"{report_type.capitalize()} Brand Intelligence Report"
    p2.font.name = "Arial"
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = text_white
    p2.space_after = Pt(20)
    p2.alignment = PP_ALIGN.LEFT
    
    p3 = tf.add_paragraph()
    p3.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
    p3.font.name = "Arial"
    p3.font.size = Pt(16)
    p3.font.color.rgb = text_muted
    p3.alignment = PP_ALIGN.LEFT

    # Slide 2: Executive Overview
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Executive Performance Overview"
    p.font.name = "Arial"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = text_white
    
    # 4 KPI Cards side-by-side
    kpis = [
        ("Reputation Score", f"{data['reputation_score']}/100", "Overall health index"),
        ("Average Star Rating", f"{data['avg_rating']} ★", "Customer satisfaction"),
        ("Reply Rate", f"{int(data['response_rate'] * 100)}%", "Review response coverage"),
        ("Positive Sentiment", f"{int(data['sentiment_score'] * 100)}%", "Ratio of positive reviews")
    ]
    
    card_width = Inches(2.7)
    card_height = Inches(3.8)
    gap = Inches(0.3)
    start_left = Inches(0.75)
    start_top = Inches(1.8)
    
    for i, (title, value, desc) in enumerate(kpis):
        left = start_left + i * (card_width + gap)
        
        # Draw background shape for card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, start_top, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = card_bg_color
        card.line.color.rgb = primary_color
        card.line.width = Pt(1.5)
        
        # Add content in the card
        tb = slide.shapes.add_textbox(left + Inches(0.1), start_top + Inches(0.3), card_width - Inches(0.2), card_height - Inches(0.6))
        ctf = tb.text_frame
        ctf.word_wrap = True
        
        # Card title
        cp1 = ctf.paragraphs[0]
        cp1.text = title.upper()
        cp1.font.name = "Arial"
        cp1.font.size = Pt(11)
        cp1.font.bold = True
        cp1.font.color.rgb = text_muted
        cp1.alignment = PP_ALIGN.CENTER
        cp1.space_after = Pt(40)
        
        # Card value
        cp2 = ctf.add_paragraph()
        cp2.text = value
        cp2.font.name = "Arial"
        cp2.font.size = Pt(36)
        cp2.font.bold = True
        cp2.font.color.rgb = text_white
        cp2.alignment = PP_ALIGN.CENTER
        cp2.space_after = Pt(40)
        
        # Card description
        cp3 = ctf.add_paragraph()
        cp3.text = desc
        cp3.font.name = "Arial"
        cp3.font.size = Pt(10)
        cp3.font.color.rgb = text_muted
        cp3.alignment = PP_ALIGN.CENTER

    # Slide 3: Sentiment & Top Topics
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Sentiment & Operational Topics"
    p.font.name = "Arial"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = text_white
    
    # Sentiment column (Left)
    left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.5))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = card_bg_color
    left_box.line.color.rgb = RGBColor(51, 65, 85) # Slate 600
    
    ltb = slide.shapes.add_textbox(Inches(0.95), Inches(2.0), Inches(5.2), Inches(4.1))
    ltf = ltb.text_frame
    ltf.word_wrap = True
    
    lp1 = ltf.paragraphs[0]
    lp1.text = "SENTIMENT DISTRIBUTION"
    lp1.font.name = "Arial"
    lp1.font.size = Pt(14)
    lp1.font.bold = True
    lp1.font.color.rgb = primary_color
    lp1.space_after = Pt(20)
    
    s_dist = data["sentiment_distribution"]
    for s_name, s_count in s_dist.items():
        lp = ltf.add_paragraph()
        lp.text = f"• {s_name.capitalize()}: {s_count} reviews"
        lp.font.name = "Arial"
        lp.font.size = Pt(14)
        lp.font.color.rgb = text_white
        lp.space_after = Pt(10)
        
    # Top Topics column (Right)
    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.5))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = card_bg_color
    right_box.line.color.rgb = RGBColor(51, 65, 85)
    
    rtb = slide.shapes.add_textbox(Inches(7.18), Inches(2.0), Inches(5.2), Inches(4.1))
    rtf = rtb.text_frame
    rtf.word_wrap = True
    
    rp1 = rtf.paragraphs[0]
    rp1.text = "TOP EXTRACTED TOPICS"
    rp1.font.name = "Arial"
    rp1.font.size = Pt(14)
    rp1.font.bold = True
    rp1.font.color.rgb = primary_color
    rp1.space_after = Pt(20)
    
    top_topics = data.get("top_topics", [])
    if not top_topics:
        rp = rtf.add_paragraph()
        rp.text = "No topics extracted yet."
        rp.font.name = "Arial"
        rp.font.size = Pt(14)
        rp.font.color.rgb = text_muted
    else:
        for topic_item in top_topics[:5]:
            topic_name = topic_item["name"]
            topic_count = topic_item["count"]
            topic_sentiment = topic_item["sentiment_score"]
            rp = rtf.add_paragraph()
            rp.text = f"• {topic_name}: {topic_count} mentions (Sentiment Index: {int(topic_sentiment * 100)}%)"
            rp.font.name = "Arial"
            rp.font.size = Pt(13)
            rp.font.color.rgb = text_white
            rp.space_after = Pt(10)

    # Slide 4: AI Recommendations & Actions
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Action Center Recommendations"
    p.font.name = "Arial"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = text_white
    
    recs = data.get("ai_recommendations", [])
    
    rec_width = Inches(11.83)
    rec_height = Inches(1.2)
    start_top = Inches(1.8)
    gap = Inches(0.25)
    
    for i, rec in enumerate(recs[:3]): # Max 3 recommendations
        top = start_top + i * (rec_height + gap)
        
        # Border box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), top, rec_width, rec_height)
        box.fill.solid()
        box.fill.fore_color.rgb = card_bg_color
        box.line.color.rgb = primary_color if rec["type"] != "warning" else RGBColor(239, 68, 68)
        box.line.width = Pt(1.5)
        
        # Text
        tb = slide.shapes.add_textbox(Inches(0.95), top + Inches(0.1), rec_width - Inches(0.4), rec_height - Inches(0.2))
        rtf = tb.text_frame
        rtf.word_wrap = True
        
        rp1 = rtf.paragraphs[0]
        rp1.text = f"{rec['title'].upper()} ({rec['type'].upper()})"
        rp1.font.name = "Arial"
        rp1.font.size = Pt(12)
        rp1.font.bold = True
        rp1.font.color.rgb = primary_color if rec["type"] != "warning" else RGBColor(239, 68, 68)
        
        rp2 = rtf.add_paragraph()
        rp2.text = rec["description"]
        rp2.font.name = "Arial"
        rp2.font.size = Pt(11)
        rp2.font.color.rgb = text_white
        rp2.space_before = Pt(4)
        
    stream = io.BytesIO()
    prs.save(stream)
    return stream.getvalue()
